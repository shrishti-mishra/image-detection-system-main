from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import os

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0

def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.7), width=Inches(8))
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = f"Image not found: {image_path}"

def build_presentation(output_path):
    prs = Presentation()
    add_title_slide(prs, "Advanced Image Detection System", "Object detection, mapping, and descriptions")

    add_bullets_slide(
        prs,
        "Overview",
        [
            "Flask web app for image uploads",
            "Annotates detections with bounding boxes",
            "Generates heatmaps and descriptions",
            "Extracts location via EXIF and keywords",
            "Map view with Google Maps or Leaflet fallback",
        ],
    )

    add_bullets_slide(
        prs,
        "Architecture",
        [
            "app.py: routes and orchestration",
            "image_processor.py: annotation and heatmap",
            "location.py: EXIF and reverse geocoding",
            "description.py: summary text and tags",
            "templates/: result rendering and map",
        ],
    )

    add_bullets_slide(
        prs,
        "Detection & Labels",
        [
            "Supports normalized and pixel bboxes",
            "Fallback detection to ensure demo visibility",
            "Labels show class and confidence",
        ],
    )

    add_bullets_slide(
        prs,
        "Mapping",
        [
            "EXIF GPS → decimal degrees conversion (fixed)",
            "Reverse geocoding via Nominatim / Google",
            "Landmark enrichment (Taj Mahal, Eiffel Tower)",
            "India fallback when no coordinates",
        ],
    )

    add_bullets_slide(
        prs,
        "How To Run",
        [
            "python -m venv venv; .\\venv\\Scripts\\activate",
            "pip install Flask requests numpy opencv-python Pillow geopy python-dotenv matplotlib",
            "python image-detection-system-main\\project\\app.py",
            "Open http://127.0.0.1:5000/",
        ],
    )

    uploads_dir = os.path.join(os.getcwd(), "image-detection-system-main", "static", "uploads")
    latest_image = None
    if os.path.isdir(uploads_dir):
        files = [os.path.join(uploads_dir, f) for f in os.listdir(uploads_dir) if f.lower().endswith((".png",".jpg",".jpeg"))]
        if files:
            latest_image = sorted(files, key=os.path.getmtime, reverse=True)[0]
    add_image_slide(prs, "Recent Upload (if available)", latest_image or "")

    prs.save(output_path)

if __name__ == "__main__":
    out = os.path.join(os.getcwd(), "image-detection-system-main", "Advanced_Image_Detection_System.pptx")
    build_presentation(out)
    print("Saved:", out)
